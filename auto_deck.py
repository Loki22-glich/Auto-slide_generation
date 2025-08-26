import os
import sys
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from ddgs import DDGS
from openai import OpenAI

# =========================
# Web Search Function
# =========================
def search_web(query, num_results=5):
    results = []
    with DDGS() as ddgs:
        for r in ddgs.text(query, max_results=num_results):
            results.append(f"{r['title']}: {r['body']}")
    return results

# =========================
# LLM Content Generator
# =========================
def generate_slide_content(topic, search_snippets):
    client = OpenAI()

    prompt = f"""
    You are an assistant that creates professional slide decks.

    Topic: {topic}

    Use the search results below + your knowledge to create a structured slide outline for a 7-slide deck:

    - Slide 1: Title
    - Slide 2: Overview
    - Slide 3–6: Key points, trends, or arguments (each slide should have a title + 3–4 bullet points)
    - Slide 7: Conclusion / Takeaways

    Make it concise, clear, and presentation-ready.

    Search Results:
    {search_snippets}
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
    )

    return response.choices[0].message.content

# =========================
# Branded PPTX Builder
# =========================
def build_pptx(structured_content, output_file="Branded_Slides.pptx", brand_name="AutoDeck AI"):
    prs = Presentation()
    bullet_layout = prs.slide_layouts[1]

    # Theme styles
    title_font_size = Pt(40)
    bullet_font_size = Pt(24)
    title_color = RGBColor(0, 51, 102)     # Dark blue
    bullet_color = RGBColor(50, 50, 50)    # Dark gray
    footer_color = RGBColor(100, 100, 100) # Light gray

    # Split slides by "Slide X:"
    slides = structured_content.split("Slide ")
    for slide in slides:
        slide = slide.strip()
        if not slide or not slide[0].isdigit():
            continue

        # Extract title & body
        parts = slide.split(":", 1)
        if len(parts) < 2:
            continue
        slide_title, slide_body = parts
        slide_body = slide_body.strip()

        slide_obj = prs.slides.add_slide(bullet_layout)
        title = slide_obj.shapes.title
        content = slide_obj.placeholders[1]

        # === Add Background Rectangle ===
        background = slide_obj.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(7.5)  # full slide rectangle
        )
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 245, 255)  # Light blue background
        background.shadow.inherit = False
        slide_obj.shapes._spTree.remove(background._element)  # send to back
        slide_obj.shapes._spTree.insert(2, background._element)

        # === Title formatting ===
        title.text = slide_body.split("\n")[0] if slide_body else f"Slide {slide_title}"
        title.text_frame.paragraphs[0].font.size = title_font_size
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color

        # === Content / bullet formatting ===
        tf = content.text_frame
        tf.clear()
        for line in slide_body.split("\n")[1:]:
            line = line.strip("-• ").strip()
            if not line:
                continue

            p = tf.add_paragraph()

            # Detect markdown-style heading (***Heading*** or ## Heading)
            if line.startswith("*") or line.startswith("#"):
                clean_text = line.replace("*", "").replace("#", "").strip()
                p.text = clean_text
                p.font.size = Pt(26)  # slightly larger than bullets
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)  # brighter blue for emphasis
                p.level = 0  # top-level (like subheading)
            else:
                p.text = line
                p.font.size = bullet_font_size
                p.font.color.rgb = bullet_color
                p.level = 1  # normal bullet


        # === Footer with brand ===
        footer = slide_obj.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(12), Inches(0.5))
        footer_tf = footer.text_frame
        footer_tf.text = f"{brand_name} | Auto-Generated Deck"
        footer_tf.paragraphs[0].font.size = Pt(12)
        footer_tf.paragraphs[0].font.color.rgb = footer_color

    prs.save(output_file)
    print(f"✅ Branded slide deck saved as {output_file}")

# =========================
# Main
# =========================
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 auto_deck_step5_branded.py 'Your Topic'")
        sys.exit(1)

    topic = sys.argv[1]
    print(f"🔎 Searching web for: {topic}")
    search_results = search_web(topic)

    print(f"📑 Generating content for: {topic}")
    structured_content = generate_slide_content(topic, "\n".join(search_results))

    print("🎨 Building Branded PowerPoint...")
    build_pptx(structured_content, f"{topic.replace(' ', '_')}_Branded_Deck.pptx")
